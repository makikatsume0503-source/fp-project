import json
from datetime import datetime
from pathlib import Path
from .base_agent import BaseAgent

SYSTEM_PROMPT = """あなたはセミナー講師とプレゼンテーションデザインの専門家です。
おうちCFO FP相談サービスのセミナー資料（パワーポイント）を作成します。

スライド設計の方針:
- 対象: 共働き・子育て世帯（30〜40代）
- 構成: タイトル→アジェンダ→本編（各トピック）→まとめ→CTA
- 1スライドに1メッセージ
- 具体的な数字・事例を活用
- 行動を促す締めくくり
- 視覚的にわかりやすい図解・箇条書き"""

OUTPUT_DIR = Path(__file__).parent.parent / "output"


class SeminarCreatorAgent(BaseAgent):
    """セミナー作成エージェント。Claude でスライド構成を生成し、python-pptx で実際のPPTXを作成する。"""

    def run(self, topic: str, slide_count: int = 15) -> dict:
        prompt = f"""以下のテーマでセミナー用スライド構成を作成してください：

テーマ: {topic}
スライド枚数: 約{slide_count}枚

以下のJSON形式で回答してください:
{{
  "title": "セミナータイトル",
  "subtitle": "サブタイトル",
  "target_audience": "想定参加者",
  "duration_minutes": 60,
  "learning_objectives": ["学習目標1", "学習目標2", "学習目標3"],
  "slides": [
    {{
      "slide_number": 1,
      "type": "title",
      "title": "スライドタイトル",
      "content": ["本文1", "本文2"],
      "speaker_notes": "講師メモ",
      "visual_suggestion": "視覚要素の提案"
    }}
  ],
  "summary": "セミナー概要"
}}"""

        raw = self.call_json(
            prompt=prompt,
            system_prompt=SYSTEM_PROMPT,
            max_tokens=12000,
        )

        data = self.extract_json(raw)
        if not data:
            return {"raw_output": raw, "summary": "JSON解析に失敗しました。"}

        pptx_path = self._build_pptx(data)
        data["pptx_path"] = str(pptx_path)
        return data

    def _build_pptx(self, data: dict) -> Path:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        NAVY    = RGBColor(0x1A, 0x2B, 0x4A)
        NAVY2   = RGBColor(0x22, 0x3A, 0x5E)
        NAVY3   = RGBColor(0x0D, 0x1A, 0x35)
        GOLD    = RGBColor(0xC8, 0xA9, 0x6E)
        GOLD2   = RGBColor(0xE8, 0xD0, 0xA0)
        WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
        LGRAY   = RGBColor(0xF7, 0xF4, 0xEF)
        MGRAY   = RGBColor(0x88, 0x88, 0x88)
        ACCENT  = RGBColor(0x4A, 0x90, 0xD9)
        TEAL    = RGBColor(0x2E, 0x8B, 0x7A)
        WARM    = RGBColor(0xC8, 0x70, 0x40)

        CARD_BG  = [NAVY, NAVY2, RGBColor(0x26, 0x40, 0x66), RGBColor(0x1E, 0x35, 0x58)]
        ACCENTS  = [GOLD, ACCENT, TEAL, WARM]

        W = Inches(13.33)
        H = Inches(7.5)

        prs = Presentation()
        prs.slide_width = W
        prs.slide_height = H
        blank = prs.slide_layouts[6]
        slides_data = data.get("slides", [])
        total = len(slides_data)

        def rect(slide, x, y, w, h, color):
            s = slide.shapes.add_shape(1, x, y, w, h)
            s.fill.solid()
            s.fill.fore_color.rgb = color
            s.line.fill.background()
            return s

        def oval(slide, x, y, w, h, color):
            s = slide.shapes.add_shape(9, x, y, w, h)
            s.fill.solid()
            s.fill.fore_color.rgb = color
            s.line.fill.background()
            return s

        def text(slide, t, x, y, w, h, size, bold=False, color=None, align=PP_ALIGN.LEFT):
            tb = slide.shapes.add_textbox(x, y, w, h)
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = align
            r = p.add_run()
            r.text = t
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color or RGBColor(0x44, 0x44, 0x44)
            return tb

        def header(slide, label, title_t, label_color=GOLD):
            rect(slide, Inches(0), Inches(0), W, Inches(1.38), NAVY)
            rect(slide, Inches(0), Inches(1.38), W, Inches(0.07), GOLD)
            text(slide, label, Inches(0.42), Inches(0.1), Inches(3), Inches(0.42),
                 10, bold=True, color=label_color)
            text(slide, title_t, Inches(0.42), Inches(0.52), Inches(12.5), Inches(0.78),
                 27, bold=True, color=WHITE)

        def footer(slide, n, tot):
            rect(slide, Inches(0), Inches(7.15), W, Inches(0.05), GOLD)
            text(slide, f"{n} / {tot}", Inches(12.2), Inches(7.2),
                 Inches(1.0), Inches(0.28), 9, color=MGRAY, align=PP_ALIGN.RIGHT)
            text(slide, "おうちCFO FP相談", Inches(0.35), Inches(7.2),
                 Inches(5), Inches(0.28), 9, color=MGRAY)

        for idx, sd in enumerate(slides_data, 1):
            slide = prs.slides.add_slide(blank)
            stype   = sd.get("type", "content")
            stitle  = sd.get("title", "")
            content = sd.get("content", [])
            notes   = sd.get("speaker_notes", "")
            acc     = ACCENTS[(idx - 1) % len(ACCENTS)]

            # ==================================================
            # タイトルスライド
            # ==================================================
            if stype == "title":
                rect(slide, Inches(0), Inches(0), W, H, NAVY3)
                # 右上大円・小円装飾
                oval(slide, Inches(8.8), Inches(-2.2), Inches(6.5), Inches(6.5), NAVY2)
                oval(slide, Inches(11.2), Inches(0.3), Inches(2.4), Inches(2.4),
                     RGBColor(0x9A, 0x78, 0x38))
                oval(slide, Inches(12.4), Inches(1.5), Inches(1.0), Inches(1.0),
                     RGBColor(0xC8, 0xA9, 0x6E))
                # 左ゴールドバー
                rect(slide, Inches(0), Inches(0), Inches(0.28), H, GOLD)
                # 左上細ライン
                rect(slide, Inches(0.28), Inches(0), W, Inches(0.06), NAVY2)
                # ドット装飾
                for di in range(4):
                    oval(slide, Inches(0.8 + di * 0.45), Inches(1.85),
                         Inches(0.18), Inches(0.18),
                         GOLD if di == 0 else RGBColor(0x55, 0x66, 0x88))
                # タグライン
                tag = data.get("target_audience", "共働き・子育て世帯向け")[:30]
                text(slide, tag, Inches(0.82), Inches(1.15), Inches(10), Inches(0.55),
                     13, color=GOLD2)
                # メインタイトル
                tb = slide.shapes.add_textbox(Inches(0.82), Inches(2.15), Inches(10.8), Inches(2.3))
                tf = tb.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                r = p.add_run()
                r.text = data.get("title", stitle)
                r.font.size = Pt(48)
                r.font.bold = True
                r.font.color.rgb = WHITE
                # 区切りライン
                rect(slide, Inches(0.82), Inches(4.6), Inches(4.2), Inches(0.09), GOLD)
                # サブタイトル
                if data.get("subtitle"):
                    text(slide, data["subtitle"], Inches(0.82), Inches(4.82),
                         Inches(11), Inches(0.8), 20, color=GOLD2)
                # 発表者
                text(slide, "ファイナンシャルプランナー　勝目麻希",
                     Inches(0.82), Inches(5.78), Inches(8), Inches(0.5),
                     14, color=RGBColor(0xBB, 0xBB, 0xBB))
                # 時間バッジ
                dur = data.get("duration_minutes", 60)
                rect(slide, Inches(10.4), Inches(6.15), Inches(2.5), Inches(0.6), GOLD)
                text(slide, f"約{dur}分", Inches(10.4), Inches(6.15),
                     Inches(2.5), Inches(0.6), 13, bold=True, color=NAVY,
                     align=PP_ALIGN.CENTER)

            # ==================================================
            # アジェンダスライド
            # ==================================================
            elif stype == "agenda":
                rect(slide, Inches(0), Inches(0), W, H, LGRAY)
                rect(slide, Inches(0), Inches(0), W, Inches(1.5), NAVY)
                rect(slide, Inches(0), Inches(1.5), W, Inches(0.07), GOLD)
                # 右上装飾円
                oval(slide, Inches(11.5), Inches(0.1), Inches(1.2), Inches(1.2), NAVY2)

                text(slide, "AGENDA", Inches(0.5), Inches(0.13), Inches(3), Inches(0.45),
                     11, bold=True, color=GOLD)
                text(slide, stitle, Inches(0.5), Inches(0.62), Inches(11.8), Inches(0.78),
                     27, bold=True, color=WHITE)

                half = len(content) // 2 + len(content) % 2
                cols = [content[:half], content[half:]]
                ccolors = [NAVY, GOLD, ACCENT, TEAL, WARM, NAVY2]

                for ci, col in enumerate(cols):
                    x = Inches(0.7 + ci * 6.35)
                    for ri, item in enumerate(col):
                        gi = ri + ci * half
                        y = Inches(1.78 + ri * 0.97)
                        cc = ccolors[gi % len(ccolors)]
                        # 背景バー（交互）
                        rect(slide, x, y + Inches(0.03),
                             Inches(5.9), Inches(0.8),
                             RGBColor(0xEB, 0xE7, 0xDF) if ri % 2 == 0 else WHITE)
                        # 番号円
                        oval(slide, x + Inches(0.05), y + Inches(0.13),
                             Inches(0.55), Inches(0.55), cc)
                        text(slide, str(gi + 1), x + Inches(0.05), y + Inches(0.1),
                             Inches(0.55), Inches(0.6), 16, bold=True, color=WHITE,
                             align=PP_ALIGN.CENTER)
                        text(slide, item, x + Inches(0.72), y + Inches(0.18),
                             Inches(5.1), Inches(0.55), 14, color=NAVY,
                             bold=(gi == 0))

                footer(slide, idx, total)

            # ==================================================
            # まとめ・CTAスライド
            # ==================================================
            elif stype in ("summary", "cta"):
                rect(slide, Inches(0), Inches(0), W, H, NAVY3)
                oval(slide, Inches(8.5), Inches(-2.5), Inches(7), Inches(7), NAVY2)
                rect(slide, Inches(0), Inches(0), Inches(0.28), H, GOLD)
                rect(slide, Inches(0.28), Inches(0), W, Inches(0.06), NAVY2)

                label = "SUMMARY" if stype == "summary" else "NEXT STEP"
                text(slide, label, Inches(0.72), Inches(0.2), Inches(4), Inches(0.48),
                     11, bold=True, color=GOLD)
                text(slide, stitle, Inches(0.72), Inches(0.72), Inches(11.5), Inches(0.82),
                     30, bold=True, color=WHITE)
                rect(slide, Inches(0.72), Inches(1.62), Inches(3.5), Inches(0.08), GOLD)

                per_row = 3 if len(content) >= 5 else 2
                for ci, item in enumerate(content):
                    row, col = ci // per_row, ci % per_row
                    cw = Inches(12.5 / per_row)
                    cx = Inches(0.42) + col * (cw + Inches(0.12))
                    cy = Inches(1.85) + row * Inches(2.4)
                    bc = CARD_BG[ci % len(CARD_BG)]
                    rect(slide, cx, cy, cw, Inches(2.15), bc)
                    # ゴールド上ボーダー
                    rect(slide, cx, cy, cw, Inches(0.09), GOLD)
                    # 番号円
                    oval(slide, cx + Inches(0.15), cy + Inches(0.2),
                         Inches(0.45), Inches(0.45), GOLD)
                    text(slide, str(ci + 1), cx + Inches(0.15), cy + Inches(0.18),
                         Inches(0.45), Inches(0.45), 12, bold=True, color=NAVY,
                         align=PP_ALIGN.CENTER)
                    tb = slide.shapes.add_textbox(
                        cx + Inches(0.15), cy + Inches(0.72), cw - Inches(0.3), Inches(1.32))
                    tf = tb.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    r = p.add_run()
                    r.text = item
                    r.font.size = Pt(14)
                    r.font.color.rgb = WHITE

                footer(slide, idx, total)

            # ==================================================
            # ステップスライド（横フロー）
            # ==================================================
            elif stype == "step":
                rect(slide, Inches(0), Inches(0), W, H, LGRAY)
                header(slide, "STEP", stitle)
                rect(slide, Inches(0), Inches(0), Inches(0.18), H, acc)

                n = len(content)
                sw = (Inches(12.6) - Inches(0.25) * (n - 1)) / n
                step_colors = [NAVY, ACCENT, TEAL, WARM, NAVY2]

                for si, item in enumerate(content):
                    sx = Inches(0.38) + si * (sw + Inches(0.25))
                    sy = Inches(1.65)
                    sc = step_colors[si % len(step_colors)]
                    rect(slide, sx, sy, sw, Inches(5.1), sc)
                    rect(slide, sx, sy, sw, Inches(0.1), GOLD if si % 2 == 0 else WHITE)
                    # ステップ番号（大）
                    text(slide, f"0{si + 1}", sx, sy + Inches(0.18), sw, Inches(0.85),
                         36, bold=True, color=GOLD if si % 2 == 0 else WHITE,
                         align=PP_ALIGN.CENTER)
                    # 区切りライン
                    rect(slide, sx + sw * 0.2, sy + Inches(1.1), sw * 0.6, Inches(0.04),
                         RGBColor(0xFF, 0xFF, 0xFF) if si % 2 == 0 else GOLD)
                    tb = slide.shapes.add_textbox(
                        sx + Inches(0.12), sy + Inches(1.25), sw - Inches(0.24), Inches(3.7))
                    tf = tb.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    r = p.add_run()
                    r.text = item
                    r.font.size = Pt(16 if n <= 4 else 13)
                    r.font.color.rgb = WHITE
                    # 矢印
                    if si < n - 1:
                        text(slide, "▶", sx + sw + Inches(0.02), sy + Inches(2.1),
                             Inches(0.21), Inches(0.55), 13, color=NAVY, align=PP_ALIGN.CENTER)

                footer(slide, idx, total)

            # ==================================================
            # ケーススタディスライド（サイドバーレイアウト）
            # ==================================================
            elif stype == "case":
                rect(slide, Inches(0), Inches(0), W, H, LGRAY)
                # 左サイドバー（ネイビー）
                rect(slide, Inches(0), Inches(0), Inches(3.9), H, NAVY)
                rect(slide, Inches(3.9), Inches(0), Inches(0.07), H, GOLD)
                # サイドバー装飾円
                oval(slide, Inches(0.5), Inches(5.0), Inches(2.8), Inches(2.8), NAVY2)

                text(slide, "CASE", Inches(0.35), Inches(0.3), Inches(3.2), Inches(0.5),
                     11, bold=True, color=GOLD)
                text(slide, f"#{idx}", Inches(0.35), Inches(0.85), Inches(3.2), Inches(1.2),
                     54, bold=True, color=WHITE)
                rect(slide, Inches(0.35), Inches(2.1), Inches(2.5), Inches(0.07), GOLD)
                text(slide, stitle, Inches(0.35), Inches(2.3), Inches(3.35), Inches(4.0),
                     18, bold=True, color=WHITE)

                # 右コンテンツ
                tb = slide.shapes.add_textbox(
                    Inches(4.25), Inches(0.55), Inches(8.8), Inches(6.6))
                tf = tb.text_frame
                tf.word_wrap = True
                for ci, item in enumerate(content):
                    p = tf.paragraphs[0] if ci == 0 else tf.add_paragraph()
                    p.space_before = Pt(14)
                    r = p.add_run()
                    r.text = f"▶  {item}"
                    r.font.size = Pt(18)
                    r.font.color.rgb = NAVY if ci % 2 == 0 else RGBColor(0x33, 0x55, 0x88)

                footer(slide, idx, total)

            # ==================================================
            # 通常コンテンツスライド（項目数でレイアウト変化）
            # ==================================================
            else:
                rect(slide, Inches(0), Inches(0), W, H, LGRAY)
                header(slide, "POINT", stitle)
                rect(slide, Inches(0), Inches(0), Inches(0.18), H, acc)

                n = len(content)

                if n == 1:
                    # 1項目: 大ハイライトボックス
                    rect(slide, Inches(0.5), Inches(1.72), Inches(12.3), Inches(5.0), NAVY)
                    rect(slide, Inches(0.5), Inches(1.72), Inches(0.14), Inches(5.0), GOLD)
                    oval(slide, Inches(10.5), Inches(2.0), Inches(2.5), Inches(2.5), NAVY2)
                    tb = slide.shapes.add_textbox(
                        Inches(0.95), Inches(2.1), Inches(11.6), Inches(4.2))
                    tf = tb.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    r = p.add_run()
                    r.text = content[0]
                    r.font.size = Pt(24)
                    r.font.color.rgb = WHITE

                elif n == 2:
                    # 2項目: 左右大ブロック
                    colors2 = [NAVY, ACCENT]
                    for ci, item in enumerate(content):
                        cx = Inches(0.45 + ci * 6.45)
                        rect(slide, cx, Inches(1.72), Inches(6.1), Inches(5.0), colors2[ci])
                        rect(slide, cx, Inches(1.72), Inches(6.1), Inches(0.1),
                             GOLD if ci == 0 else WHITE)
                        text(slide, str(ci + 1), cx + Inches(0.2), Inches(1.9),
                             Inches(0.85), Inches(0.95), 42, bold=True,
                             color=GOLD if ci == 0 else WHITE)
                        rect(slide, cx + Inches(0.2), Inches(2.95), Inches(2.5), Inches(0.06),
                             GOLD if ci == 0 else RGBColor(0xCC, 0xCC, 0xFF))
                        tb = slide.shapes.add_textbox(
                            cx + Inches(0.2), Inches(3.1), Inches(5.7), Inches(3.4))
                        tf = tb.text_frame
                        tf.word_wrap = True
                        p = tf.paragraphs[0]
                        r = p.add_run()
                        r.text = item
                        r.font.size = Pt(18)
                        r.font.color.rgb = WHITE

                elif n == 3:
                    # 3項目: 3カラムカード
                    c3 = [NAVY, RGBColor(0x26, 0x40, 0x66), NAVY2]
                    for ci, item in enumerate(content):
                        cx = Inches(0.45 + ci * 4.28)
                        cy = Inches(1.72)
                        rect(slide, cx, cy, Inches(4.0), Inches(5.0), c3[ci])
                        rect(slide, cx, cy, Inches(4.0), Inches(0.12), GOLD)
                        oval(slide, cx + Inches(1.65), cy + Inches(0.32),
                             Inches(0.7), Inches(0.7), GOLD)
                        text(slide, str(ci + 1), cx + Inches(1.65), cy + Inches(0.3),
                             Inches(0.7), Inches(0.7), 20, bold=True, color=NAVY,
                             align=PP_ALIGN.CENTER)
                        tb = slide.shapes.add_textbox(
                            cx + Inches(0.2), cy + Inches(1.2), Inches(3.6), Inches(3.6))
                        tf = tb.text_frame
                        tf.word_wrap = True
                        p = tf.paragraphs[0]
                        p.alignment = PP_ALIGN.CENTER
                        r = p.add_run()
                        r.text = item
                        r.font.size = Pt(16)
                        r.font.color.rgb = WHITE

                elif n == 4:
                    # 4項目: 2×2グリッド
                    c4 = [NAVY, RGBColor(0x26, 0x40, 0x66), NAVY2, RGBColor(0x1E, 0x35, 0x58)]
                    for ci, item in enumerate(content):
                        row, col = ci // 2, ci % 2
                        cx = Inches(0.45 + col * 6.42)
                        cy = Inches(1.67) + row * Inches(2.68)
                        rect(slide, cx, cy, Inches(6.1), Inches(2.45), c4[ci])
                        # アクセントボーダー
                        rect(slide, cx, cy, Inches(6.1), Inches(0.1),
                             GOLD if col == 0 else ACCENT)
                        text(slide, str(ci + 1), cx + Inches(0.22), cy + Inches(0.18),
                             Inches(0.6), Inches(0.65), 26, bold=True,
                             color=GOLD if col == 0 else ACCENT)
                        tb = slide.shapes.add_textbox(
                            cx + Inches(0.88), cy + Inches(0.18), Inches(5.05), Inches(2.1))
                        tf = tb.text_frame
                        tf.word_wrap = True
                        p = tf.paragraphs[0]
                        r = p.add_run()
                        r.text = item
                        r.font.size = Pt(17)
                        r.font.color.rgb = WHITE

                else:
                    # 5項目以上: 2カラムリスト
                    half = len(content) // 2 + len(content) % 2
                    for ci, col_items in enumerate([content[:half], content[half:]]):
                        cx = Inches(0.42 + ci * 6.42)
                        item_h = Inches(5.55 / half)
                        for li, item in enumerate(col_items):
                            cy = Inches(1.65) + li * item_h
                            bc = NAVY if li % 2 == 0 else NAVY2
                            rect(slide, cx, cy + Inches(0.05),
                                 Inches(6.1), item_h - Inches(0.1), bc)
                            rect(slide, cx, cy + Inches(0.05),
                                 Inches(0.1), item_h - Inches(0.1),
                                 GOLD if ci == 0 else ACCENT)
                            tb = slide.shapes.add_textbox(
                                cx + Inches(0.22), cy + Inches(0.18),
                                Inches(5.7), item_h - Inches(0.28))
                            tf = tb.text_frame
                            tf.word_wrap = True
                            p = tf.paragraphs[0]
                            r = p.add_run()
                            r.text = item
                            r.font.size = Pt(15)
                            r.font.color.rgb = WHITE

                footer(slide, idx, total)

            if notes:
                slide.notes_slide.notes_text_frame.text = notes

        OUTPUT_DIR.mkdir(exist_ok=True)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_title = data.get("title", "seminar")[:30].replace("/", "_").replace(" ", "_")
        filename = f"seminar_{safe_title}_{timestamp}.pptx"
        path = OUTPUT_DIR / filename
        prs.save(str(path))
        return path
